import os
import json
import pandas as pd
import requests
import pdfplumber
from pptx import Presentation
from docx import Document
from urllib.parse import urlparse, urljoin
import xml.etree.ElementTree as ET
from sentence_transformers import SentenceTransformer
import psycopg2
from pgvector.psycopg2 import register_vector
import warnings
import sys
import logging

warnings.filterwarnings("ignore", category=UserWarning, module='torch.distributed.elastic')
logging.getLogger("pdfminer").setLevel(logging.ERROR)

# === Configuration ===
DB_CONFIG = {
    "dbname": "market_db",
    "user": "postgres",
    "password": "1329",
    "host": "localhost",
    "port": "5432"
}
TABLE_NAME = "your_table"
BATCH_SIZE = 32
MAX_PAGES = 1000

conn = psycopg2.connect(**DB_CONFIG)
register_vector(conn)

embedder = SentenceTransformer("all-MiniLM-L6-v2")

# === Utility: Meaningful Split ===
def split_text(text, max_chunk_size=300):
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks = []
    current_chunk = ""
    for para in paragraphs:
        if len(current_chunk) + len(para.split()) <= max_chunk_size:
            current_chunk += " " + para
        else:
            chunks.append(current_chunk.strip())
            current_chunk = para
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

# === Extractors with Table-Aware Text Extraction ===
def extract_text_from_pdf(file_path, max_pages=MAX_PAGES):
    logging.getLogger("pdfminer").setLevel(logging.ERROR)  # Suppress pdfminer warnings
    text = ""

    try:
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            print(f"📄 PDF has {total_pages} pages. Extracting up to {max_pages}.")

            for i, page in enumerate(pdf.pages):
                if i >= max_pages:
                    print(f"⚠️ Reached max page limit of {max_pages}, stopping PDF extraction.")
                    break

                try:
                    # Extract paragraph text
                    page_text = page.extract_text() or ""

                    # Extract tables safely
                    tables = page.extract_tables() or []
                    table_texts = []

                    for table in tables:
                        for row in table:
                            # Safely convert each cell to string, replacing None with empty string
                            safe_row = [str(cell).strip() if cell is not None else "" for cell in row]
                            table_texts.append("\t".join(safe_row))

                    table_text = "\n".join(table_texts)

                    # Combine paragraph text and table text
                    combined_text = (page_text + "\n" + table_text).strip()
                    if combined_text:
                        text += combined_text + "\n"

                except Exception as page_err:
                    print(f"⚠️ Skipping page {i + 1} due to error: {page_err}")

                if (i + 1) % 50 == 0:
                    print(f"✅ Processed {i + 1} pages...")

        print(f"📝 Done. Extracted {len(text)} characters from {min(total_pages, max_pages)} pages.")
        return text

    except Exception as e:
        print(f"❌ PDF extraction error: {e}")
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"❌ DOCX extraction error: {e}")
        return ""

def extract_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"❌ TXT extraction error: {e}")
        return ""

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"❌ PPTX extraction error: {e}")
        return ""

def extract_text_from_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        return "\n".join(df.astype(str).apply(lambda row: " ".join(row), axis=1))
    except Exception as e:
        print(f"❌ CSV extraction error: {e}")
        return ""

def extract_text_from_json(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return json.dumps(json.load(f), indent=2)
    except Exception as e:
        print(f"❌ JSON extraction error: {e}")
        return ""

def extract_text_from_url(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return response.text
    except Exception as e:
        print(f"❌ URL extraction error: {e}")
        return ""

# === Ingestion ===
def ingest_file(file_path_or_url, website_id):
    ext = os.path.splitext(file_path_or_url)[1].lower()
    text = ""

    try:
        if file_path_or_url.startswith("http") and not file_path_or_url.endswith(".xml"):
            text = extract_text_from_url(file_path_or_url)
        elif ext == ".pdf":
            text = extract_text_from_pdf(file_path_or_url)
        elif ext == ".docx":
            text = extract_text_from_docx(file_path_or_url)
        elif ext == ".txt":
            text = extract_text_from_txt(file_path_or_url)
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_path_or_url)
        elif ext == ".csv":
            text = extract_text_from_csv(file_path_or_url)
        elif ext == ".json":
            text = extract_text_from_json(file_path_or_url)
        else:
            print(f"⚠️ Unsupported file type: {ext}")
            return
    except Exception as e:
        print(f"❌ Extraction failed: {e}")
        return

    if not text.strip():
        print("⚠️ No text extracted.")
        return

    chunks = split_text(text)
    print(f"✅ {len(chunks)} meaningful chunks prepared.")

    with conn.cursor() as cur:
        for i in range(0, len(chunks), BATCH_SIZE):
            batch = chunks[i:i+BATCH_SIZE]
            try:
                embeddings = embedder.encode(batch).tolist()
                for chunk, emb in zip(batch, embeddings):
                    cur.execute(f"""
                        INSERT INTO {TABLE_NAME} (content, embedding, website_id)
                        VALUES (%s, %s, %s)
                    """, (chunk, emb, website_id))
                print(f"✅ Batch {i//BATCH_SIZE + 1} saved.")
            except Exception as e:
                print(f"❌ Embedding/Insert failed: {e}")
        conn.commit()

def get_urls_from_sitemap(sitemap_url):
    try:
        response = requests.get(sitemap_url, timeout=10)
        root = ET.fromstring(response.content)
        return [elem.text for elem in root.iter() if elem.tag.endswith("loc")]
    except:
        return []

def auto_ingest(input_str, website_id):
    if input_str.startswith("http") and input_str.endswith(".xml"):
        urls = get_urls_from_sitemap(input_str)
        for url in urls:
            ingest_file(url, website_id)
    else:
        ingest_file(input_str, website_id)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python ingest_to_db.py <file_path_or_url> <website_id>")
        sys.exit(1)

    auto_ingest(sys.argv[1], int(sys.argv[2]))
